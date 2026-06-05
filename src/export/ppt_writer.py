from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches


def _add_table_to_slide(
    slide: Slide, df: pd.DataFrame, left: float, top: float, width: float, height: float
) -> None:
    rows = len(df) + 1
    cols = len(df.columns)
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table

    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = str(col_name)

    for row_idx, row in enumerate(df.itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)


def export_powerpoint(
    insights: list[str],
    summary_tables: dict[str, pd.DataFrame],
    charts: dict[str, Path],
    output_path: Path,
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Retail Performance Report"
    slide.placeholders[1].text = "Auto-generated with Python"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Insights"
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for idx, insight in enumerate(insights):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = insight

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Revenue by Region"
    slide.shapes.add_picture(
        str(charts["region_chart"]), Inches(0.8), Inches(1.4), width=Inches(8.0)
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Monthly Revenue Trend"
    slide.shapes.add_picture(
        str(charts["monthly_chart"]), Inches(0.8), Inches(1.4), width=Inches(8.0)
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top Regions"
    _add_table_to_slide(
        slide, summary_tables["by_region"].head(5), left=0.6, top=1.4, width=8.5, height=3.2
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top Products"
    _add_table_to_slide(
        slide, summary_tables["by_product"].head(5), left=0.6, top=1.4, width=8.5, height=3.2
    )

    prs.save(output_path)
